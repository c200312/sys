"""
AI PPT 生成模块 - 简化版
只提供核心的 PPT 生成功能，用于课程资源管理
"""
import os
import io
import base64
import uuid
import aiohttp
import uvicorn
from typing import Optional, List
from fastapi import FastAPI, UploadFile, File, HTTPException, Header, Request, Form
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from fastapi.responses import Response
from pydantic import BaseModel
from dotenv import load_dotenv
from openai import OpenAI
from pptx import Presentation
from pptx.util import Inches
from PIL import Image

# 加载后端主目录的环境变量
load_dotenv(os.path.join(os.path.dirname(__file__), "..", ".env"))

# 获取当前模块目录
MODULE_DIR = os.path.dirname(__file__)

app = FastAPI(title="AI PPT Generator", version="1.0.0")

# 允许跨域
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# 图片跨域处理
@app.middleware("http")
async def add_cors_headers_to_images(request: Request, call_next):
    response = await call_next(request)
    if request.url.path.startswith("/images/"):
        response.headers["Access-Control-Allow-Origin"] = "*"
        response.headers["Cache-Control"] = "no-cache"
    return response

# 静态文件服务
IMAGES_DIR = os.path.join(MODULE_DIR, "storage", "images")
os.makedirs(IMAGES_DIR, exist_ok=True)
app.mount("/images", StaticFiles(directory=IMAGES_DIR), name="images")

# ==================== 数据模型 ====================

class ApiKeyRequest(BaseModel):
    api_key: str

class GeneratePPTRequest(BaseModel):
    prompt: str
    page_count: int = 5
    style_template: Optional[str] = None

class SlideContent(BaseModel):
    index: int
    image_url: str
    prompt: str

class GeneratePPTResponse(BaseModel):
    success: bool
    slides: List[SlideContent] = []
    session_id: Optional[str] = None
    error: Optional[str] = None

class ExportPPTRequest(BaseModel):
    session_id: str
    image_urls: List[str]
    title: str = "AI Generated PPT"

# ==================== 核心功能 ====================

# 默认风格模板
DEFAULT_STYLE = """
A modern Tech/Internet company presentation slide.
Style: Modern SaaS aesthetic, clean UI, sleek vector art, soft shadows (glassmorphism).
Background: Clean LIGHT background (white or very light grey) with SUBTLE tech accents.
Content: Minimalist infographics, rounded cards, sans-serif typography.
Avoid: Old-school academic look, heavy dark borders, realistic photos, cluttered text.

【重要】所有幻灯片中的文字必须使用简体中文，包括标题、正文、图表标签等。
"""

def get_openai_client(api_key: str = None):
    """获取 OpenAI 客户端"""
    key = api_key or os.getenv("OPENROUTER_API_KEY")
    if not key:
        raise HTTPException(400, "API Key not configured")
    return OpenAI(
        base_url=os.getenv("OPENROUTER_BASE_URL", "https://openrouter.ai/api/v1"),
        api_key=key
    )

def plan_ppt_outline(client: OpenAI, prompt: str, page_count: int) -> List[dict]:
    """使用 LLM 规划 PPT 大纲"""
    system_prompt = f"""
    You are a presentation designer. Plan a {page_count}-slide presentation.

    Rules:
    1. Slide 1: Title/Cover slide
    2. Middle slides: Main content
    3. Last slide: Summary/Thank you
    4. 必须使用简体中文输出所有内容，包括标题和visual_prompt
    5. Each slide needs a visual_prompt for image generation
    6. visual_prompt中必须明确要求"所有文字使用简体中文"

    Output JSON format:
    {{
        "slides": [
            {{"index": 0, "title": "...", "visual_prompt": "..."}},
            ...
        ]
    }}
    """

    try:
        response = client.chat.completions.create(
            model="google/gemini-2.0-flash-001",
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": prompt}
            ],
            response_format={"type": "json_object"}
        )
        import json
        result = json.loads(response.choices[0].message.content)
        return result.get("slides", [])
    except Exception as e:
        print(f"Plan error: {e}")
        # 返回默认结构
        return [
            {"index": i, "title": f"Slide {i+1}", "visual_prompt": f"{prompt} - slide {i+1}"}
            for i in range(page_count)
        ]

def generate_slide_image(client: OpenAI, visual_prompt: str, style: str) -> str:
    """生成单张幻灯片图片"""
    chinese_requirement = "【重要要求】幻灯片中所有文字必须使用简体中文，包括标题、正文、图表标签、说明文字等。"
    full_prompt = f"{style}\n\n{chinese_requirement}\n\n**SLIDE CONTENT**: {visual_prompt}"

    try:
        response = client.chat.completions.create(
            model="google/gemini-3-pro-image-preview",
            messages=[{"role": "user", "content": full_prompt}],
            extra_body={"modalities": ["image", "text"]}
        )

        message = response.choices[0].message
        if hasattr(message, 'images') and message.images:
            return message.images[0]["image_url"]["url"]
        return None
    except Exception as e:
        print(f"Image generation error: {e}")
        return None

async def save_image(image_data: str, session_id: str) -> str:
    """保存图片到本地"""
    session_dir = os.path.join(IMAGES_DIR, session_id)
    os.makedirs(session_dir, exist_ok=True)

    filename = f"{uuid.uuid4()}.png"
    filepath = os.path.join(session_dir, filename)

    # Base64 数据
    if image_data.startswith("data:"):
        header, encoded = image_data.split(",", 1)
        data = base64.b64decode(encoded)
        with open(filepath, "wb") as f:
            f.write(data)
    # HTTP URL
    elif image_data.startswith("http"):
        async with aiohttp.ClientSession() as session:
            async with session.get(image_data) as resp:
                if resp.status == 200:
                    data = await resp.read()
                    with open(filepath, "wb") as f:
                        f.write(data)
                else:
                    return None
    # 纯 Base64
    else:
        try:
            data = base64.b64decode(image_data)
            with open(filepath, "wb") as f:
                f.write(data)
        except:
            return None

    return f"/images/{session_id}/{filename}"

# ==================== API 接口 ====================

@app.get("/")
def health_check():
    return {"status": "running", "service": "AI PPT Generator"}

@app.get("/api/key")
def get_api_key_status():
    """获取 API Key 状态"""
    key = os.getenv("OPENROUTER_API_KEY", "")
    if key:
        return {"has_key": True, "key_preview": key[:10] + "..."}
    return {"has_key": False}

@app.post("/api/key/test")
def test_api_key(req: ApiKeyRequest):
    """测试 API Key"""
    try:
        client = get_openai_client(req.api_key)
        client.models.list()
        return {"status": "success", "message": "API Key is valid"}
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.post("/ppt/generate", response_model=GeneratePPTResponse)
async def generate_ppt(
    req: GeneratePPTRequest,
    x_api_key: Optional[str] = Header(None, alias="x-api-key")
):
    """
    生成 PPT
    - prompt: 主题/提示词
    - page_count: 页数 (默认5)
    - style_template: 可选的风格模板
    """
    try:
        client = get_openai_client(x_api_key)
        style = req.style_template or DEFAULT_STYLE
        session_id = str(uuid.uuid4())

        # 1. 规划大纲
        print(f"Planning PPT: {req.prompt}")
        outline = plan_ppt_outline(client, req.prompt, req.page_count)

        # 2. 生成每张幻灯片
        slides = []
        for slide_info in outline:
            visual_prompt = slide_info.get("visual_prompt", slide_info.get("title", ""))
            print(f"Generating slide {slide_info['index']}: {visual_prompt[:50]}...")

            # 生成图片
            image_url = generate_slide_image(client, visual_prompt, style)
            if not image_url:
                continue

            # 保存到本地
            local_path = await save_image(image_url, session_id)
            if local_path:
                slides.append(SlideContent(
                    index=slide_info["index"],
                    image_url=local_path,
                    prompt=visual_prompt
                ))

        if not slides:
            return GeneratePPTResponse(success=False, error="Failed to generate slides")

        return GeneratePPTResponse(success=True, slides=slides, session_id=session_id)

    except HTTPException as e:
        return GeneratePPTResponse(success=False, error=e.detail)
    except Exception as e:
        print(f"Generate PPT error: {e}")
        return GeneratePPTResponse(success=False, error=str(e))

@app.post("/ppt/generate_single")
async def generate_single_slide(
    prompt: str = Form(...),
    style_template: Optional[str] = Form(None),
    x_api_key: Optional[str] = Header(None, alias="x-api-key")
):
    """生成单张幻灯片"""
    try:
        client = get_openai_client(x_api_key)
        style = style_template or DEFAULT_STYLE
        session_id = str(uuid.uuid4())

        image_url = generate_slide_image(client, prompt, style)
        if not image_url:
            raise HTTPException(500, "Image generation failed")

        local_path = await save_image(image_url, session_id)
        if not local_path:
            raise HTTPException(500, "Failed to save image")

        return {"success": True, "image_url": local_path}

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, str(e))


@app.post("/ppt/export")
async def export_ppt(req: ExportPPTRequest):
    """
    将幻灯片图片导出为 PPTX 文件
    返回 base64 编码的 PPTX 文件内容
    """
    try:
        # 创建 PPT
        prs = Presentation()
        # 设置幻灯片大小为 16:9
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # 获取空白布局
        blank_layout = prs.slide_layouts[6]  # 空白布局

        for image_url in req.image_urls:
            # 获取图片的完整路径
            if image_url.startswith("/images/"):
                image_path = os.path.join(IMAGES_DIR, image_url.replace("/images/", ""))
            else:
                image_path = image_url

            if not os.path.exists(image_path):
                print(f"Image not found: {image_path}")
                continue

            # 添加幻灯片
            slide = prs.slides.add_slide(blank_layout)

            # 获取图片尺寸并计算适配大小
            with Image.open(image_path) as img:
                img_width, img_height = img.size

            # 计算缩放比例，使图片适应幻灯片
            slide_width = prs.slide_width
            slide_height = prs.slide_height

            width_ratio = slide_width / Inches(img_width / 96)  # 假设 96 DPI
            height_ratio = slide_height / Inches(img_height / 96)

            # 使用较小的比例以确保图片完全显示
            if width_ratio < height_ratio:
                width = slide_width
                height = Inches(img_height / 96) * width_ratio
            else:
                height = slide_height
                width = Inches(img_width / 96) * height_ratio

            # 居中放置
            left = (slide_width - width) / 2
            top = (slide_height - height) / 2

            # 添加图片到幻灯片
            slide.shapes.add_picture(image_path, left, top, width, height)

        # 保存到内存
        pptx_buffer = io.BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)

        # 转为 base64
        pptx_base64 = base64.b64encode(pptx_buffer.read()).decode('utf-8')

        return {
            "success": True,
            "pptx_base64": f"data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{pptx_base64}",
            "filename": f"{req.title}.pptx"
        }

    except Exception as e:
        print(f"Export PPT error: {e}")
        return {"success": False, "error": str(e)}


def start():
    """启动服务"""
    port = int(os.getenv("AIPPT_PORT", 8002))
    print(f"🚀 Starting AI PPT Generator on http://localhost:{port}")
    uvicorn.run("backend.aippt.main:app", host="0.0.0.0", port=port, reload=True)


if __name__ == "__main__":
    start()
