"""
文件解析器：支持多种文档格式

解析策略：
- PDF: 使用 Qwen-VL 逐页解析图片（视觉识别）
- DOCX/DOC: 先转为 PDF，再通过 Qwen-VL 逐页解析图片
- PPT/PPTX: 先转为 PDF，再通过 Qwen-VL 逐页解析图片
- TXT/MD: 直接读取文本

图片处理：
- 从 PDF 中提取嵌入的图片
- 上传到 MinIO 存储
- VLM 生成图片描述
- Markdown 中插入 ![描述](url) 占位符
"""
import io
import os
import re
import uuid
import base64
import asyncio
import logging
import tempfile
from typing import Optional, List, Tuple
from dataclasses import dataclass

import httpx
from pypdf import PdfReader
import fitz  # PyMuPDF

logger = logging.getLogger(__name__)

# API 配置
DASHSCOPE_API_BASE = "https://dashscope.aliyuncs.com/compatible-mode/v1"

# 图片描述提示词
IMAGE_DESCRIPTION_PROMPT = """请描述这张图片的内容。要求：
1. 如果是图表/流程图/架构图，描述其结构和关键信息
2. 如果是截图/界面图，描述界面元素和功能
3. 如果是照片/插图，描述主要内容和场景
4. 用简洁的中文描述，控制在50-150字
5. 只输出描述内容，不要添加"这是"等开头"""

# 页面解析提示词
PAGE_PARSE_PROMPT = """请将这张文档页面的内容提取为纯Markdown格式。

要求：
1. 保留标题层级（#、##、###）
2. 保留列表结构（-、*、1.）
3. 保留表格结构
4. 如果页面中有图片，用 [图片] 标记位置
5. 直接输出Markdown内容，不要用 ```markdown ``` 代码块包裹
6. 不要添加任何说明文字"""


def _clean_markdown_wrapper(content: str) -> str:
    """
    清理 VLM 返回内容中的 Markdown 代码块包裹

    有些 VLM 会用 ```markdown ... ``` 包裹返回内容
    """
    import re

    # 匹配 ```markdown ... ``` 或 ``` ... ```
    pattern = r'^```(?:markdown)?\s*\n?(.*?)\n?```\s*$'
    match = re.match(pattern, content.strip(), re.DOTALL)
    if match:
        return match.group(1).strip()

    # 也处理开头有 ```markdown 但结尾没有 ``` 的情况
    if content.strip().startswith('```markdown'):
        content = content.strip()[len('```markdown'):].strip()
        if content.endswith('```'):
            content = content[:-3].strip()

    return content


@dataclass
class ExtractedImage:
    """提取的图片信息"""
    image_bytes: bytes      # 图片二进制数据
    page_num: int           # 所在页码
    image_index: int        # 页内图片索引
    width: int              # 图片宽度
    height: int             # 图片高度
    minio_url: str = ""     # MinIO 存储 URL
    description: str = ""   # VLM 生成的描述


class FileParser:
    """
    文件解析器

    支持的格式：
    - 文本文件：txt, md（直接读取）
    - PDF 文档：pdf（Qwen-VL 视觉识别）
    - Word 文档：doc, docx（转 PDF 后用 Qwen-VL 解析）
    - PPT 文档：ppt, pptx（转 PDF 后用 Qwen-VL 解析）
    """

    # 支持的文件扩展名
    SUPPORTED_TEXT = {'txt', 'md'}
    SUPPORTED_DOCX = {'docx', 'doc'}  # Word 文档：转 PDF 后用 Qwen-VL
    SUPPORTED_QWEN_VL = {'pdf'}  # PDF 使用 Qwen-VL
    SUPPORTED_PPT = {'ppt', 'pptx'}

    @classmethod
    async def parse(
        cls,
        filename: str,
        content_bytes: bytes,
        knowledge_id: str = ""
    ) -> str:
        """
        解析文件内容

        Args:
            filename: 文件名（用于判断类型）
            content_bytes: 文件二进制内容
            knowledge_id: 知识库 ID（用于图片存储路径）

        Returns:
            解析后的 Markdown 文本内容
        """
        ext = filename.lower().split('.')[-1] if '.' in filename else ''
        file_size = len(content_bytes)

        logger.info(f"[PARSER] ========== 开始解析文件 ==========")
        logger.info(f"[PARSER] 文件名: {filename}")
        logger.info(f"[PARSER] 扩展名: {ext}")
        logger.info(f"[PARSER] 文件大小: {file_size} bytes ({file_size/1024:.2f} KB)")
        logger.info(f"[PARSER] 知识库ID: {knowledge_id or '(未指定)'}")

        try:
            if ext in cls.SUPPORTED_TEXT:
                logger.info(f"[PARSER] 解析方式: 直接读取文本")
                result = cls._parse_text(content_bytes)
            elif ext in cls.SUPPORTED_QWEN_VL:
                logger.info(f"[PARSER] 解析方式: Qwen-VL 视觉识别 (PDF)")
                result = await cls._parse_pdf_direct(content_bytes, knowledge_id)
            elif ext in cls.SUPPORTED_DOCX:
                logger.info(f"[PARSER] 解析方式: LibreOffice + Qwen-VL (Word)")
                result = await cls._parse_docx(filename, content_bytes, knowledge_id)
            elif ext in cls.SUPPORTED_PPT:
                logger.info(f"[PARSER] 解析方式: LibreOffice + Qwen-VL (PPT)")
                result = await cls._parse_ppt(filename, content_bytes, knowledge_id)
            else:
                logger.info(f"[PARSER] 解析方式: 尝试作为文本解析 (未知类型)")
                result = cls._parse_text(content_bytes)

            logger.info(f"[PARSER] 解析完成: {len(result)} 字符")
            logger.info(f"[PARSER] ----- 解析内容预览(前500字) -----")
            preview_lines = result[:500].split('\n')
            for line in preview_lines:
                if line.strip():
                    logger.info(f"[PARSER]   {line}")
            if len(result) > 500:
                logger.info(f"[PARSER]   ... (省略 {len(result) - 500} 字)")
            logger.info(f"[PARSER] ========== 解析结束 ==========")
            return result

        except Exception as e:
            logger.error(f"[PARSER] 文件解析错误 [{filename}]: {e}", exc_info=True)
            # 降级到本地解析
            logger.info(f"[PARSER] 降级到本地解析...")
            return await cls._fallback_parse(filename, content_bytes)

    @classmethod
    def _parse_text(cls, content_bytes: bytes) -> str:
        """解析纯文本文件"""
        return content_bytes.decode('utf-8', errors='ignore')

    @classmethod
    async def _parse_pdf_direct(cls, content_bytes: bytes, knowledge_id: str = "") -> str:
        """
        直接使用 Qwen-VL 解析 PDF 文件

        流程：
        1. PDF 转图片（每页）
        2. 使用 Qwen-VL 解析每页图片
        3. 提取嵌入图片并生成描述
        4. 合并所有页面的解析结果
        """
        logger.info(f"[PARSER] ----- PDF直接解析开始 -----")

        api_key = os.getenv("DASHSCOPE_API_KEY")
        if not api_key:
            logger.warning("[PARSER] 阿里云 API Key 未配置，使用本地解析")
            return cls._local_parse_pdf(content_bytes)

        logger.info(f"[PARSER] DASHSCOPE_API_KEY 已配置")

        try:
            result = await cls._parse_pdf_with_qwen_vl(content_bytes, api_key, knowledge_id)
            if result:
                return result
            else:
                logger.warning("[PARSER] Qwen-VL 解析结果为空，降级使用本地解析")
                return cls._local_parse_pdf(content_bytes)
        except Exception as e:
            logger.error(f"[PARSER] PDF解析失败: {e}", exc_info=True)
            return cls._local_parse_pdf(content_bytes)

    @classmethod
    async def _parse_docx(cls, filename: str, content_bytes: bytes, knowledge_id: str = "") -> str:
        """
        解析 Word 文件

        流程：
        1. Word 转 PDF（使用 LibreOffice）
        2. PDF 转图片（每页）
        3. 使用 Qwen-VL 解析每页图片
        4. 提取嵌入图片并生成描述
        5. 合并所有页面的解析结果
        """
        logger.info(f"[PARSER] ----- Word解析开始 -----")

        api_key = os.getenv("DASHSCOPE_API_KEY")
        if not api_key:
            logger.warning("[PARSER] 阿里云 API Key 未配置，使用本地解析")
            return cls._local_parse_docx(content_bytes)

        logger.info(f"[PARSER] DASHSCOPE_API_KEY 已配置")

        try:
            # 将 Word 保存为临时文件
            ext = filename.lower().split('.')[-1] if '.' in filename else 'docx'
            suffix = f'.{ext}'
            with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp_doc:
                tmp_doc.write(content_bytes)
                tmp_doc_path = tmp_doc.name

            logger.info(f"[PARSER] Word临时文件: {tmp_doc_path}")

            try:
                # 使用 LibreOffice 转换 Word 为 PDF
                logger.info(f"[PARSER] 开始转换Word为PDF...")
                pdf_bytes = await cls._convert_office_to_pdf(tmp_doc_path)

                if pdf_bytes:
                    logger.info(f"[PARSER] Word转PDF成功，PDF大小: {len(pdf_bytes)} bytes")
                    # PDF 转图片并解析
                    logger.info(f"[PARSER] 开始使用Qwen-VL解析PDF...")
                    return await cls._parse_pdf_with_qwen_vl(pdf_bytes, api_key, knowledge_id)
                else:
                    # 降级：使用本地解析
                    logger.warning(f"[PARSER] Word转PDF失败，降级使用本地解析")
                    return cls._local_parse_docx(content_bytes)

            finally:
                # 清理临时文件
                if os.path.exists(tmp_doc_path):
                    os.unlink(tmp_doc_path)
                    logger.info(f"[PARSER] 已清理临时文件")

        except Exception as e:
            logger.error(f"[PARSER] Word解析失败: {e}", exc_info=True)
            return cls._local_parse_docx(content_bytes)

    @classmethod
    async def _parse_ppt(cls, filename: str, content_bytes: bytes, knowledge_id: str = "") -> str:
        """
        解析 PPT 文件

        流程：
        1. PPT 转 PDF（使用 LibreOffice）
        2. PDF 转图片（每页）
        3. 使用 Qwen-VL 解析每页图片
        4. 提取嵌入图片并生成描述
        5. 合并所有页面的解析结果
        """
        logger.info(f"[PARSER] ----- PPT解析开始 -----")

        api_key = os.getenv("DASHSCOPE_API_KEY")
        if not api_key:
            logger.warning("[PARSER] 阿里云 API Key 未配置，使用本地解析")
            return await cls._fallback_parse(filename, content_bytes)

        logger.info(f"[PARSER] DASHSCOPE_API_KEY 已配置")

        try:
            # 将 PPT 保存为临时文件
            suffix = '.pptx' if filename.lower().endswith('.pptx') else '.ppt'
            with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp_ppt:
                tmp_ppt.write(content_bytes)
                tmp_ppt_path = tmp_ppt.name

            logger.info(f"[PARSER] PPT临时文件: {tmp_ppt_path}")

            try:
                # 使用 LibreOffice 转换 PPT 为 PDF
                logger.info(f"[PARSER] 开始转换PPT为PDF...")
                pdf_bytes = await cls._convert_office_to_pdf(tmp_ppt_path)

                if pdf_bytes:
                    logger.info(f"[PARSER] PPT转PDF成功，PDF大小: {len(pdf_bytes)} bytes")
                    # PDF 转图片并解析
                    logger.info(f"[PARSER] 开始使用Qwen-VL解析PDF...")
                    return await cls._parse_pdf_with_qwen_vl(pdf_bytes, api_key, knowledge_id)
                else:
                    # 降级：直接解析 PPT 文本
                    logger.warning(f"[PARSER] PPT转PDF失败，降级使用本地解析")
                    return await cls._fallback_parse(filename, content_bytes)

            finally:
                # 清理临时文件
                if os.path.exists(tmp_ppt_path):
                    os.unlink(tmp_ppt_path)
                    logger.info(f"[PARSER] 已清理临时文件")

        except Exception as e:
            logger.error(f"[PARSER] PPT解析失败: {e}", exc_info=True)
            return await cls._fallback_parse(filename, content_bytes)

    @classmethod
    async def _convert_office_to_pdf(cls, file_path: str) -> Optional[bytes]:
        """
        将 Office 文档（Word/PPT）转换为 PDF

        使用 LibreOffice 进行转换
        """
        try:
            import subprocess
            import shutil

            # 创建临时输出目录
            output_dir = tempfile.mkdtemp()
            logger.info(f"[PARSER] LibreOffice输出目录: {output_dir}")

            try:
                # 尝试使用 LibreOffice
                # Windows: 检查常见安装路径, Linux/Mac: libreoffice
                soffice_cmd = None
                if os.name == "nt":
                    # Windows 常见安装路径
                    possible_paths = [
                        r"C:\Program Files\LibreOffice\program\soffice.exe",
                        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
                    ]
                    for path in possible_paths:
                        if os.path.exists(path):
                            soffice_cmd = path
                            logger.info(f"[PARSER] 找到LibreOffice: {path}")
                            break
                    if not soffice_cmd:
                        # 尝试 PATH 中的 soffice
                        soffice_cmd = "soffice"
                        logger.info(f"[PARSER] 使用PATH中的soffice")
                else:
                    soffice_cmd = "libreoffice"
                    logger.info(f"[PARSER] 使用libreoffice命令")

                logger.info(f"[PARSER] 执行LibreOffice转换命令...")
                result = subprocess.run(
                    [
                        soffice_cmd,
                        "--headless",
                        "--convert-to", "pdf",
                        "--outdir", output_dir,
                        file_path
                    ],
                    capture_output=True,
                    timeout=60
                )

                logger.info(f"[PARSER] LibreOffice返回码: {result.returncode}")
                if result.stdout:
                    logger.info(f"[PARSER] LibreOffice stdout: {result.stdout.decode('utf-8', errors='ignore')}")
                if result.stderr:
                    logger.warning(f"[PARSER] LibreOffice stderr: {result.stderr.decode('utf-8', errors='ignore')}")

                if result.returncode == 0:
                    # 查找生成的 PDF
                    pdf_filename = os.path.splitext(os.path.basename(file_path))[0] + ".pdf"
                    pdf_path = os.path.join(output_dir, pdf_filename)
                    logger.info(f"[PARSER] 预期PDF路径: {pdf_path}")

                    if os.path.exists(pdf_path):
                        logger.info(f"[PARSER] PDF文件存在，读取中...")
                        with open(pdf_path, 'rb') as f:
                            pdf_data = f.read()
                        logger.info(f"[PARSER] PDF读取成功，大小: {len(pdf_data)} bytes")
                        return pdf_data
                    else:
                        logger.warning(f"[PARSER] PDF文件不存在: {pdf_path}")
                        # 列出目录内容
                        dir_contents = os.listdir(output_dir)
                        logger.info(f"[PARSER] 输出目录内容: {dir_contents}")
                else:
                    logger.warning(f"[PARSER] LibreOffice转换失败，返回码: {result.returncode}")

            except FileNotFoundError as e:
                logger.warning(f"[PARSER] LibreOffice未安装或找不到: {e}")
            except subprocess.TimeoutExpired:
                logger.warning("[PARSER] Office转PDF超时(60秒)")
            finally:
                # 清理临时目录
                shutil.rmtree(output_dir, ignore_errors=True)
                logger.info(f"[PARSER] 已清理临时输出目录")

        except Exception as e:
            logger.error(f"[PARSER] PPT转PDF失败: {e}", exc_info=True)

        return None

    # 并行处理的最大并发数
    MAX_CONCURRENT_VLM_CALLS = 5

    @classmethod
    async def _parse_pdf_with_qwen_vl(
        cls,
        pdf_bytes: bytes,
        api_key: str,
        knowledge_id: str = ""
    ) -> str:
        """
        使用 Qwen-VL 并行解析 PDF 的每一页，并提取嵌入图片

        流程：
        1. PDF 逐页渲染为图片，并行发送给 VLM 解析文本
        2. 提取 PDF 中嵌入的图片（去重）
        3. 并行上传图片到 MinIO，VLM 生成描述
        4. 在 Markdown 中插入图片占位符

        Args:
            pdf_bytes: PDF 二进制内容
            api_key: DashScope API Key
            knowledge_id: 知识库 ID（用于图片存储路径）
        """
        try:
            doc = fitz.open(stream=pdf_bytes, filetype="pdf")
            total_pages = len(doc)
            logger.info(f"[PARSER] PDF共 {total_pages} 页，开始并行解析（最大并发: {cls.MAX_CONCURRENT_VLM_CALLS}）...")

            # 第一步：提取所有嵌入图片（全局去重，记录首次出现的页码）
            logger.info(f"[PARSER] ----- 开始提取嵌入图片 -----")
            all_images, image_page_map = cls._extract_all_images(doc)
            logger.info(f"[PARSER] 共提取到 {len(all_images)} 张唯一图片")

            # 第二步：预先渲染所有页面为图片（同步操作，但很快）
            logger.info(f"[PARSER] ----- 开始渲染页面图片 -----")
            page_images_data = []
            for page_num in range(total_pages):
                page = doc[page_num]
                mat = fitz.Matrix(2, 2)
                pix = page.get_pixmap(matrix=mat)
                img_bytes = pix.tobytes("png")
                img_base64 = base64.b64encode(img_bytes).decode('utf-8')
                page_images_data.append((page_num, img_base64, len(img_bytes)))
            logger.info(f"[PARSER] 页面渲染完成，共 {len(page_images_data)} 页")

            doc.close()

            # 第三步：并行处理图片上传和描述生成
            if all_images:
                logger.info(f"[PARSER] ----- 开始并行处理嵌入图片 -----")
                all_images = await cls._process_extracted_images_parallel(
                    all_images, api_key, knowledge_id
                )

            # 第四步：并行解析所有页面内容
            logger.info(f"[PARSER] ----- 开始并行解析页面内容 -----")
            semaphore = asyncio.Semaphore(cls.MAX_CONCURRENT_VLM_CALLS)

            async def parse_single_page(page_num: int, img_base64: str, img_size: int) -> tuple:
                """解析单个页面"""
                async with semaphore:
                    logger.info(f"[PARSER] 开始解析第 {page_num + 1}/{total_pages} 页 ({img_size} bytes)...")
                    async with httpx.AsyncClient(timeout=60.0) as client:
                        try:
                            response = await client.post(
                                f"{DASHSCOPE_API_BASE}/chat/completions",
                                headers={
                                    "Authorization": f"Bearer {api_key}",
                                    "Content-Type": "application/json"
                                },
                                json={
                                    "model": "qwen-vl-plus",
                                    "messages": [
                                        {
                                            "role": "user",
                                            "content": [
                                                {
                                                    "type": "text",
                                                    "text": PAGE_PARSE_PROMPT
                                                },
                                                {
                                                    "type": "image_url",
                                                    "image_url": {
                                                        "url": f"data:image/png;base64,{img_base64}"
                                                    }
                                                }
                                            ]
                                        }
                                    ],
                                    "max_tokens": 4096
                                }
                            )

                            if response.status_code == 200:
                                result = response.json()
                                page_content = result.get("choices", [{}])[0].get("message", {}).get("content", "")
                                if page_content:
                                    # 清理可能的 markdown 代码块包裹
                                    page_content = _clean_markdown_wrapper(page_content)
                                    logger.info(f"[PARSER] 第 {page_num + 1} 页解析成功，{len(page_content)} 字符")
                                else:
                                    page_content = "[内容为空]"
                                    logger.warning(f"[PARSER] 第 {page_num + 1} 页解析结果为空")
                            else:
                                page_content = "[解析失败]"
                                logger.warning(f"[PARSER] 第 {page_num + 1} 页 API 失败: {response.status_code}")

                            return (page_num, page_content)

                        except Exception as e:
                            logger.error(f"[PARSER] 第 {page_num + 1} 页解析异常: {e}")
                            return (page_num, "[解析失败]")

            # 并行执行所有页面解析
            tasks = [
                parse_single_page(page_num, img_base64, img_size)
                for page_num, img_base64, img_size in page_images_data
            ]
            results = await asyncio.gather(*tasks)

            # 按页码排序结果
            results.sort(key=lambda x: x[0])

            # 第五步：组装最终内容，插入图片占位符
            # 注意：不使用 ## 页码标记，避免干扰后续的语义分块
            all_content = []
            for page_num, page_content in results:
                # 替换 [图片] 占位符为实际的图片 Markdown
                page_images = [img for img in all_images if img.page_num == page_num]
                if page_images:
                    page_content = cls._insert_image_placeholders(page_content, page_images)
                all_content.append(page_content)

            # 用空行连接各页内容，保持文档连续性
            final_content = "\n\n".join(all_content)
            logger.info(f"[PARSER] Qwen-VL并行解析完成，总内容长度: {len(final_content)} 字符")
            return final_content

        except Exception as e:
            logger.error(f"[PARSER] Qwen-VL解析PDF失败: {e}", exc_info=True)
            return ""

    @classmethod
    def _extract_all_images(cls, doc) -> Tuple[List[ExtractedImage], dict]:
        """
        从整个 PDF 文档提取所有唯一的嵌入图片

        通过 xref 去重，记录每张图片首次出现的页码

        Args:
            doc: PyMuPDF 文档对象

        Returns:
            (图片列表, {xref: 首次出现页码})
        """
        images = []
        seen_xrefs = {}  # xref -> 首次出现的页码
        image_index = 0

        for page_num in range(len(doc)):
            page = doc[page_num]
            image_list = page.get_images(full=True)

            for img_info in image_list:
                xref = img_info[0]

                # 跳过已处理的图片
                if xref in seen_xrefs:
                    continue

                try:
                    base_image = doc.extract_image(xref)
                    if not base_image:
                        continue

                    image_bytes = base_image["image"]
                    width = base_image.get("width", 0)
                    height = base_image.get("height", 0)

                    # 过滤太小的图片（可能是图标或装饰）
                    if width < 50 or height < 50:
                        logger.debug(f"[PARSER] 跳过小图片 xref={xref}: {width}x{height}")
                        continue

                    # 转换为 PNG 格式（统一格式）
                    if base_image.get("ext") != "png":
                        try:
                            from PIL import Image
                            img = Image.open(io.BytesIO(image_bytes))
                            png_buffer = io.BytesIO()
                            img.save(png_buffer, format="PNG")
                            image_bytes = png_buffer.getvalue()
                        except Exception as e:
                            logger.warning(f"[PARSER] 图片格式转换失败: {e}")

                    seen_xrefs[xref] = page_num
                    images.append(ExtractedImage(
                        image_bytes=image_bytes,
                        page_num=page_num,  # 首次出现的页码
                        image_index=image_index,
                        width=width,
                        height=height
                    ))
                    image_index += 1
                    logger.info(f"[PARSER] 提取图片 xref={xref} 首次出现在第 {page_num + 1} 页, {width}x{height}")

                except Exception as e:
                    logger.warning(f"[PARSER] 提取图片失败 xref={xref}: {e}")

        return images, seen_xrefs

    @classmethod
    async def _process_extracted_images(
        cls,
        images: List[ExtractedImage],
        api_key: str,
        knowledge_id: str
    ) -> List[ExtractedImage]:
        """
        处理提取的图片：上传到 MinIO 并生成 VLM 描述（串行版本，已废弃）

        Args:
            images: 提取的图片列表
            api_key: DashScope API Key
            knowledge_id: 知识库 ID

        Returns:
            处理后的图片列表（包含 URL 和描述）
        """
        # 使用并行版本
        return await cls._process_extracted_images_parallel(images, api_key, knowledge_id)

    @classmethod
    async def _process_extracted_images_parallel(
        cls,
        images: List[ExtractedImage],
        api_key: str,
        knowledge_id: str
    ) -> List[ExtractedImage]:
        """
        并行处理提取的图片：上传到 MinIO 并生成 VLM 描述

        Args:
            images: 提取的图片列表
            api_key: DashScope API Key
            knowledge_id: 知识库 ID

        Returns:
            处理后的图片列表（包含 URL 和描述）
        """
        if not images:
            return images

        try:
            from app.services.minio_service import get_minio_service
            minio_service = get_minio_service()
        except Exception as e:
            logger.warning(f"[PARSER] MinIO 服务不可用，跳过图片上传: {e}")
            return images

        semaphore = asyncio.Semaphore(cls.MAX_CONCURRENT_VLM_CALLS)

        async def process_single_image(img: ExtractedImage, index: int) -> None:
            """处理单个图片"""
            async with semaphore:
                try:
                    # 1. 上传到 MinIO（同步操作，很快）
                    img_uuid = str(uuid.uuid4())[:8]
                    filename = f"page{img.page_num + 1}_img{img.image_index + 1}_{img_uuid}.png"
                    object_name = minio_service.upload_file(
                        file_data=img.image_bytes,
                        filename=filename,
                        content_type="image/png",
                        prefix="rag_images/",
                        folder_id=knowledge_id or str(uuid.uuid4())
                    )
                    img.minio_url = object_name
                    logger.info(f"[PARSER] 图片 {index+1}/{len(images)} 上传成功: {object_name}")

                    # 2. 调用 VLM 生成描述
                    img_base64 = base64.b64encode(img.image_bytes).decode('utf-8')
                    async with httpx.AsyncClient(timeout=30.0) as client:
                        response = await client.post(
                            f"{DASHSCOPE_API_BASE}/chat/completions",
                            headers={
                                "Authorization": f"Bearer {api_key}",
                                "Content-Type": "application/json"
                            },
                            json={
                                "model": "qwen-vl-plus",
                                "messages": [
                                    {
                                        "role": "user",
                                        "content": [
                                            {"type": "text", "text": IMAGE_DESCRIPTION_PROMPT},
                                            {
                                                "type": "image_url",
                                                "image_url": {
                                                    "url": f"data:image/png;base64,{img_base64}"
                                                }
                                            }
                                        ]
                                    }
                                ],
                                "max_tokens": 500
                            }
                        )

                        if response.status_code == 200:
                            result = response.json()
                            description = result.get("choices", [{}])[0].get("message", {}).get("content", "")
                            img.description = description.strip() or "图片"
                            logger.info(f"[PARSER] 图片 {index+1} 描述: {img.description[:50]}...")
                        else:
                            img.description = "图片"
                            logger.warning(f"[PARSER] 图片 {index+1} 描述生成失败: {response.status_code}")

                except Exception as e:
                    logger.warning(f"[PARSER] 处理图片 {index+1} 失败: {e}")
                    img.description = "图片"

        # 并行处理所有图片
        tasks = [process_single_image(img, i) for i, img in enumerate(images)]
        await asyncio.gather(*tasks)

        return images

    @classmethod
    def _insert_image_placeholders(cls, content: str, images: List[ExtractedImage]) -> str:
        """
        将 [图片] 占位符替换为实际的 Markdown 图片语法

        如果 VLM 没有识别到 [图片]，则在内容末尾追加图片

        Args:
            content: 页面内容
            images: 该页的图片列表

        Returns:
            替换后的内容
        """
        if not images:
            return content

        # 统计 [图片] 占位符数量
        placeholder_pattern = r'\[图片\]|\[图\]|\[image\]|\[IMAGE\]'
        placeholders = re.findall(placeholder_pattern, content, re.IGNORECASE)
        placeholder_count = len(placeholders)

        # 生成图片 Markdown
        image_markdowns = []
        for img in images:
            if img.minio_url:
                # 使用 MinIO object_name 作为图片路径
                # 前端需要通过 API 获取预签名 URL
                md = f"\n\n![{img.description}](minio://{img.minio_url})\n\n"
            else:
                md = f"\n\n[图片: {img.description}]\n\n"
            image_markdowns.append(md)

        if placeholder_count > 0:
            # 依次替换占位符
            result = content
            for i, md in enumerate(image_markdowns):
                if i < placeholder_count:
                    result = re.sub(placeholder_pattern, md, result, count=1)
            # 剩余图片追加到末尾
            for md in image_markdowns[placeholder_count:]:
                result += md
            return result
        else:
            # 没有占位符，追加到内容末尾
            return content + "\n\n" + "".join(image_markdowns)

    @classmethod
    async def _fallback_parse(cls, filename: str, content_bytes: bytes) -> str:
        """
        降级本地解析方案

        当 API 不可用时使用本地库解析
        """
        ext = filename.lower().split('.')[-1] if '.' in filename else ''
        logger.info(f"[PARSER] ----- 使用本地降级解析 -----")
        logger.info(f"[PARSER] 文件类型: {ext}")

        try:
            if ext == 'pdf':
                logger.info(f"[PARSER] 本地解析PDF...")
                result = cls._local_parse_pdf(content_bytes)
            elif ext == 'docx':
                logger.info(f"[PARSER] 本地解析DOCX...")
                result = cls._local_parse_docx(content_bytes)
            elif ext in ('ppt', 'pptx'):
                logger.info(f"[PARSER] 本地解析PPT/PPTX...")
                result = cls._local_parse_pptx(content_bytes)
            else:
                logger.info(f"[PARSER] 作为文本解析...")
                result = cls._parse_text(content_bytes)

            logger.info(f"[PARSER] 本地解析完成，内容长度: {len(result)} 字符")
            if result:
                preview = result[:200].replace('\n', ' ')
                logger.info(f"[PARSER] 内容预览: {preview}...")
            return result
        except Exception as e:
            logger.error(f"[PARSER] 本地解析失败 [{filename}]: {e}", exc_info=True)
            return ""

    @classmethod
    def _local_parse_pdf(cls, content_bytes: bytes) -> str:
        """本地解析 PDF"""
        try:
            pdf_reader = PdfReader(io.BytesIO(content_bytes))
            text_parts = []
            for page in pdf_reader.pages:
                text = page.extract_text()
                if text:
                    text_parts.append(text)
            return '\n\n'.join(text_parts)
        except Exception as e:
            logger.error(f"本地 PDF 解析失败: {e}")
            return ""

    @classmethod
    def _local_parse_docx(cls, content_bytes: bytes) -> str:
        """本地解析 Word 文档"""
        try:
            from docx import Document as DocxDocument
            doc = DocxDocument(io.BytesIO(content_bytes))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            return '\n\n'.join(paragraphs)
        except Exception as e:
            logger.error(f"本地 DOCX 解析失败: {e}")
            return ""

    @classmethod
    def _local_parse_pptx(cls, content_bytes: bytes) -> str:
        """本地解析 PPT 文档（仅提取文本）"""
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(content_bytes))
            all_text = []

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = [f"## 第 {slide_num} 页"]
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                all_text.append('\n\n'.join(slide_text))

            return '\n\n---\n\n'.join(all_text)
        except Exception as e:
            logger.error(f"本地 PPTX 解析失败: {e}")
            return ""

    @classmethod
    def is_supported(cls, filename: str) -> bool:
        """
        检查文件是否支持解析

        Args:
            filename: 文件名

        Returns:
            是否支持
        """
        ext = filename.lower().split('.')[-1] if '.' in filename else ''
        all_supported = cls.SUPPORTED_TEXT | cls.SUPPORTED_DOCX | cls.SUPPORTED_QWEN_VL | cls.SUPPORTED_PPT
        return ext in all_supported


# 便捷函数（异步版本）
async def parse_file_content(
    filename: str,
    content_bytes: bytes,
    knowledge_id: str = ""
) -> str:
    """解析文件内容（便捷函数）"""
    return await FileParser.parse(filename, content_bytes, knowledge_id)


# 同步包装器（兼容旧代码）
def parse_file_content_sync(
    filename: str,
    content_bytes: bytes,
    knowledge_id: str = ""
) -> str:
    """解析文件内容（同步版本，用于兼容）"""
    try:
        loop = asyncio.get_event_loop()
        if loop.is_running():
            # 如果已经在异步上下文中，创建新任务
            import concurrent.futures
            with concurrent.futures.ThreadPoolExecutor() as executor:
                future = executor.submit(
                    asyncio.run,
                    FileParser.parse(filename, content_bytes, knowledge_id)
                )
                return future.result()
        else:
            return loop.run_until_complete(
                FileParser.parse(filename, content_bytes, knowledge_id)
            )
    except RuntimeError:
        return asyncio.run(FileParser.parse(filename, content_bytes, knowledge_id))
